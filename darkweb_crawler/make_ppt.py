"""
Darkweb Analyzer 발표용 PPT 생성 스크립트
실행: python3 make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ──────────────────────────────────────────
BG_DARK     = RGBColor(0x07, 0x08, 0x16)
BG_CARD     = RGBColor(0x0f, 0x17, 0x2a)
PRIMARY     = RGBColor(0x63, 0x66, 0xf1)
PURPLE      = RGBColor(0x8b, 0x5c, 0xf6)
CYAN        = RGBColor(0x22, 0xd3, 0xee)
GREEN       = RGBColor(0x10, 0xb9, 0x81)
RED         = RGBColor(0xef, 0x44, 0x44)
AMBER       = RGBColor(0xf5, 0x9e, 0x0b)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LIGHT       = RGBColor(0xe2, 0xe8, 0xf0)
MUTED       = RGBColor(0x64, 0x74, 0x8b)
INDIGO_SOFT = RGBColor(0xc7, 0xd2, 0xfe)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=BG_DARK):
    bg_shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    return bg_shape


def box(slide, l, t, w, h, fill=BG_CARD, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def accent_bar(slide, l, t, h=0.04, color=PRIMARY):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(0.04), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return txb


def label(slide, text, l, t, color=PRIMARY):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(len(text)*0.12+0.3), Inches(0.28))
    s.fill.solid()
    s.fill.fore_color.rgb = BG_CARD
    s.line.color.rgb = color
    s.line.width = Pt(1)
    txb = slide.shapes.add_textbox(Inches(l+0.06), Inches(t+0.02), Inches(len(text)*0.12+0.18), Inches(0.24))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Pretendard"


def logo_box(slide, l=0.4, t=0.22):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(0.4), Inches(0.4))
    s.fill.solid()
    s.fill.fore_color.rgb = PRIMARY
    s.line.fill.background()
    txt(slide, "D", l+0.07, t+0.04, 0.26, 0.3, size=14, bold=True, color=WHITE)


def slide_header(slide, title, subtitle=None, tag=None):
    logo_box(slide)
    txt(slide, "Darkweb Analyzer", 0.9, 0.25, 3, 0.3, size=11, bold=True, color=MUTED)
    # divider
    d = slide.shapes.add_shape(1, Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.012))
    d.fill.solid(); d.fill.fore_color.rgb = RGBColor(0x1e, 0x29, 0x3b); d.line.fill.background()
    if tag:
        label(slide, tag, 0.4, 0.9)
    txt(slide, title, 0.4, 1.25, 12, 0.7, size=30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 2.0, 12, 0.45, size=14, color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ═══════════════════════════════════════════════════════════
s1 = add_slide(); bg(s1)

# 중앙 히어로 텍스트
txt(s1, "DARKWEB INTELLIGENCE", 1.5, 1.6, 10, 0.4, size=11, bold=True,
    color=PRIMARY, align=PP_ALIGN.CENTER)

title_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.6))
title_box.word_wrap = False
tf = title_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Darkweb Domain Analyzer"
run.font.size = Pt(46)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Pretendard"

txt(s1, ".onion 도메인 종합 분석 시스템", 1.5, 3.75, 10, 0.5, size=18,
    color=INDIGO_SOFT, align=PP_ALIGN.CENTER)
txt(s1, "Tor 크롤링 · CoDA 범죄 분류 · BART 유형 분류 · GPT-4o-mini 위험도 분석",
    1.5, 4.25, 10, 0.4, size=13, color=MUTED, align=PP_ALIGN.CENTER)

# 하단 구분선
d = s1.shapes.add_shape(1, Inches(4.5), Inches(4.85), Inches(4.3), Inches(0.03))
d.fill.solid(); d.fill.fore_color.rgb = PRIMARY; d.line.fill.background()

txt(s1, "보안 연구 목적  |  DarkBERT 기반 AI 분류", 1.5, 5.1, 10, 0.35,
    size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ═══════════════════════════════════════════════════════════
s2 = add_slide(); bg(s2)
slide_header(s2, "목차", tag="AGENDA")

items = [
    ("01", "프로젝트 개요",        "개발 배경 및 목적"),
    ("02", "시스템 아키텍처",       "전체 구조 및 데이터 흐름"),
    ("03", "핵심 기능 — CoDA 분류", "DarkBERT 기반 범죄 카테고리 분류"),
    ("04", "핵심 기능 — 유형 분류", "BART zero-shot 사이트 유형 분류"),
    ("05", "핵심 기능 — AI 분석",   "GPT-4o-mini 위험도 평가"),
    ("06", "보고서 및 UI",          "웹 인터페이스 및 결과 시각화"),
    ("07", "기술 스택",             "사용 기술 및 모델"),
]

for i, (num, title, sub) in enumerate(items):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    l = 0.5 + col * 6.5
    t = 2.65 + row * 0.9
    b = box(s2, l, t, 6.0, 0.75, fill=BG_CARD)
    accent_bar(s2, l, t, h=0.75, color=PRIMARY)
    txt(s2, num, l+0.18, t+0.1, 0.6, 0.3, size=11, bold=True, color=PRIMARY)
    txt(s2, title, l+0.18, t+0.28, 3.5, 0.3, size=13, bold=True, color=WHITE)
    txt(s2, sub, l+0.18, t+0.5, 5.5, 0.25, size=10, color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — 프로젝트 개요
# ═══════════════════════════════════════════════════════════
s3 = add_slide(); bg(s3)
slide_header(s3, "프로젝트 개요", "다크웹 .onion 도메인 종합 분석 자동화", tag="01 OVERVIEW")

# 배경
box(s3, 0.4, 2.55, 5.9, 2.1, fill=BG_CARD)
accent_bar(s3, 0.4, 2.55, h=2.1, color=CYAN)
txt(s3, "개발 배경", 0.6, 2.65, 5.5, 0.35, size=12, bold=True, color=CYAN)
bg_text = (
    "다크웹은 일반 검색엔진에 색인되지 않는 .onion 도메인으로 구성된 "
    "익명 네트워크로, 마약·무기·해킹 서비스 등 불법 거래의 온상이 됩니다. "
    "기존 수동 분석 방식의 한계를 극복하고 AI 기반 자동 분류 시스템을 구축했습니다."
)
txt(s3, bg_text, 0.6, 3.05, 5.5, 1.5, size=12, color=LIGHT)

# 목적
box(s3, 6.55, 2.55, 6.3, 2.1, fill=BG_CARD)
accent_bar(s3, 6.55, 2.55, h=2.1, color=GREEN)
txt(s3, "개발 목적", 6.75, 2.65, 5.9, 0.35, size=12, bold=True, color=GREEN)
goals = [
    "• .onion 도메인 접근성 자동 확인 (Tor)",
    "• DarkBERT AI 모델로 범죄 콘텐츠 분류",
    "• BART 모델로 사이트 유형 자동 판별",
    "• GPT-4o-mini 기반 위험도 평가",
    "• 결과를 HTML 보고서로 시각화",
]
for j, g in enumerate(goals):
    txt(s3, g, 6.75, 3.0 + j*0.33, 5.9, 0.3, size=12, color=LIGHT)

# 핵심 수치
stats = [("93%", "CoDA 분류 정확도"), ("9개", "범죄 카테고리"), ("3개", "AI 모델 통합")]
for i, (val, lbl) in enumerate(stats):
    l = 0.4 + i * 4.3
    box(s3, l, 4.85, 3.9, 1.15, fill=BG_CARD)
    txt(s3, val, l, 4.9, 3.9, 0.65, size=34, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    txt(s3, lbl, l, 5.52, 3.9, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — 시스템 아키텍처
# ═══════════════════════════════════════════════════════════
s4 = add_slide(); bg(s4)
slide_header(s4, "시스템 아키텍처", "전체 데이터 흐름 및 컴포넌트 구성", tag="02 ARCHITECTURE")

# 4단계 메인 플로우
flow4 = [
    ("①", "Domain Input",
     "web/app.py · agent.py",
     "사용자가 분석할\n도메인을 입력합니다.",
     PRIMARY),
    ("②", "Remote Server",
     "server/app.py (포트 5001)",
     "Tor 네트워크를 통해\n원격 서버에서\n데이터를 수집합니다.",
     CYAN),
    ("③", "Local Analysis\nModule",
     "analyzers/",
     "수집된 데이터를 기반으로\nAI/NLP 분석 알고리즘을\n통해 특성을 분석합니다.",
     PURPLE),
    ("④", "Report\nGeneration",
     "agent_report_generator.py",
     "분석 결과를 시각화하여\n보고서를 생성하고\n사용자에게 제공합니다.",
     AMBER),
]

box_w = 2.9
box_h = 3.4
start_l = 0.5
gap = 0.36

for i, (num, title, sub, desc, color) in enumerate(flow4):
    l = start_l + i * (box_w + gap)

    # 카드 박스
    b = box(s4, l, 2.2, box_w, box_h, fill=BG_CARD)

    # 상단 컬러 바
    cb = s4.shapes.add_shape(1, Inches(l), Inches(2.2), Inches(box_w), Inches(0.07))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()

    # 번호 뱃지
    nb = s4.shapes.add_shape(1, Inches(l + box_w/2 - 0.22), Inches(2.32), Inches(0.44), Inches(0.44))
    nb.fill.solid(); nb.fill.fore_color.rgb = color; nb.line.fill.background()
    txt(s4, num, l + box_w/2 - 0.22, 2.33, 0.44, 0.34,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 제목
    txt(s4, title, l, 2.85, box_w, 0.55,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 서브타이틀 (파일명)
    txt(s4, sub, l, 3.42, box_w, 0.35,
        size=9, bold=False, color=color, align=PP_ALIGN.CENTER)

    # 설명
    txt(s4, desc, l + 0.15, 3.82, box_w - 0.3, 1.6,
        size=11, color=LIGHT, align=PP_ALIGN.CENTER)

    # 화살표
    if i < len(flow4) - 1:
        ax = l + box_w + 0.04
        txt(s4, "▶", ax, 3.6, gap - 0.04, 0.45,
            size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

# 분석 모듈 세부 항목
txt(s4, "③ 분석 모듈 세부 구성", 0.4, 5.82, 12.5, 0.32, size=11, bold=True, color=MUTED)
modules = [
    ("CoDA 분류",    "DarkBERT + LogisticRegression", RED),
    ("유형 분류",    "BART zero-shot classifier",     GREEN),
    ("AI 분석",      "GPT-4o-mini (OpenAI API)",      CYAN),
    ("HTML 정제",    "BeautifulSoup4 텍스트 추출",    PURPLE),
]
for i, (name, desc, color) in enumerate(modules):
    l = 0.4 + i * 3.22
    b = box(s4, l, 6.18, 3.05, 0.6, fill=BG_CARD)
    accent_bar(s4, l, 6.18, h=0.6, color=color)
    txt(s4, name, l+0.18, 6.24, 1.3, 0.28, size=11, bold=True, color=color)
    txt(s4, desc, l+0.18, 6.5,  2.7, 0.24, size=9,  color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — CoDA 범죄 분류
# ═══════════════════════════════════════════════════════════
s5 = add_slide(); bg(s5)
slide_header(s5, "CoDA 범죄 카테고리 분류", "DarkBERT 임베딩 + LogisticRegression", tag="03 CODA")

# 파이프라인
pipe = [
    (".onion HTML 수집", CYAN),
    ("HTML 정제\n(텍스트 추출)", PURPLE),
    ("DarkBERT 임베딩\n(768차원 벡터)", PRIMARY),
    ("LogisticRegression\n(카테고리 분류)", AMBER),
    ("결과 출력\nDrugs 87%...", GREEN),
]
for i, (step, color) in enumerate(pipe):
    l = 0.35 + i * 2.56
    b = box(s5, l, 2.4, 2.3, 1.0, fill=BG_CARD)
    cb = s5.shapes.add_shape(1, Inches(l), Inches(2.4), Inches(2.3), Inches(0.06))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s5, step, l, 2.5, 2.3, 0.85, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipe)-1:
        txt(s5, "→", l+2.3, 2.75, 0.26, 0.3, size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# 카테고리 + F1
cats = [
    ("Gambling", "98%", GREEN), ("Porn",      "93%", RED),
    ("Drugs",    "92%", RED),   ("Arms",      "92%", AMBER),
    ("Electronic","91%",CYAN),  ("Violence",  "91%", RED),
    ("Financial", "87%", AMBER),("Crypto",    "82%", PRIMARY),
    ("Hacking",  "81%", PURPLE),
]
txt(s5, "9개 카테고리 F1 정확도", 0.4, 3.65, 12, 0.35, size=13, bold=True, color=MUTED)
for i, (cat, f1, color) in enumerate(cats):
    col = i % 5
    row = i // 5
    l = 0.4 + col * 2.5
    t = 4.05 + row * 0.8
    b = box(s5, l, t, 2.3, 0.68, fill=BG_CARD)
    txt(s5, cat, l+0.12, t+0.06, 1.5, 0.3, size=12, bold=True, color=color)
    txt(s5, f1,  l+0.12, t+0.34, 1.5, 0.28, size=13, bold=True, color=WHITE)

txt(s5, "전체 정확도 93%  |  검증 데이터 1,063개  |  학습 데이터 7,081개",
    0.4, 6.85, 12.5, 0.3, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — 사이트 유형 분류
# ═══════════════════════════════════════════════════════════
s6 = add_slide(); bg(s6)
slide_header(s6, "사이트 유형 분류", "BART zero-shot 분류 모델 (facebook/bart-large-mnli)", tag="04 CATEGORY")

# 설명
box(s6, 0.4, 2.55, 7.5, 2.5, fill=BG_CARD)
accent_bar(s6, 0.4, 2.55, h=2.5, color=GREEN)
txt(s6, "BART Zero-Shot 분류란?", 0.6, 2.65, 7.1, 0.4, size=13, bold=True, color=GREEN)
desc = (
    "사전 학습 없이 자연어 라벨만으로 분류합니다.\n"
    "HTML 텍스트와 각 카테고리 라벨의 자연어 연관성을\n"
    "확률로 계산하여 가장 적합한 유형을 판별합니다.\n\n"
    "• 별도 학습 데이터 불필요\n"
    "• 신규 카테고리 추가 용이\n"
    "• 로그인 페이지 자동 감지"
)
txt(s6, desc, 0.6, 3.1, 7.0, 1.85, size=12, color=LIGHT)

# 카테고리 목록
cats6 = [
    "온라인 마켓플레이스", "포럼 / 토론", "블로그",
    "뉴스 / 미디어",       "소셜 네트워크", "통신 / 메시지",
    "기술 문서 / 위키",    "제품·서비스 홍보", "로그인 필수",
    "마켓 + 포럼 혼합",    "개인 블로그",     "분류 불가",
]
txt(s6, "분류 카테고리 (13종)", 8.1, 2.55, 5, 0.4, size=13, bold=True, color=MUTED)
for i, c in enumerate(cats6):
    col = i % 2
    row = i // 2
    l = 8.1 + col * 2.55
    t = 3.0 + row * 0.52
    b = box(s6, l, t, 2.35, 0.4, fill=BG_CARD)
    accent_bar(s6, l, t, h=0.4, color=GREEN)
    txt(s6, c, l+0.18, t+0.08, 2.1, 0.28, size=10, color=LIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — AI 위험도 분석
# ═══════════════════════════════════════════════════════════
s7 = add_slide(); bg(s7)
slide_header(s7, "AI 위험도 분석", "GPT-4o-mini 기반 자연어 사이트 분석", tag="05 AI ANALYSIS")

# 출력 항목들
outputs = [
    ("사이트 유형",   "마켓플레이스, 포럼, 블로그 등",        CYAN),
    ("주요 언어",     "영어, 한국어, 러시아어 등",             PURPLE),
    ("목적",         "사이트의 주요 목적 자연어 설명",          PRIMARY),
    ("요약",         "콘텐츠 전반 요약",                       GREEN),
    ("위험도",       "낮음 / 중간 / 높음 / 매우높음",          AMBER),
    ("위험도 근거",   "판단 이유 자연어 설명",                  RED),
    ("주목할 특징",   "특이사항, 주의 필요 요소 목록",          INDIGO_SOFT),
]
txt(s7, "분석 출력 항목", 0.4, 2.55, 12, 0.35, size=13, bold=True, color=MUTED)
for i, (name, desc, color) in enumerate(outputs):
    col = i % 2
    row = i // 2
    l = 0.4 + col * 6.5
    t = 3.0 + row * 0.88
    b = box(s7, l, t, 6.1, 0.76, fill=BG_CARD)
    accent_bar(s7, l, t, h=0.76, color=color)
    txt(s7, name, l+0.18, t+0.08, 2.5, 0.3, size=12, bold=True, color=color)
    txt(s7, desc, l+0.18, t+0.4,  5.7, 0.3, size=11, color=LIGHT)

txt(s7, "* OPENAI_API_KEY 설정 시 활성화 / 미설정 시 해당 섹션 스킵",
    0.4, 6.9, 12.5, 0.3, size=10, color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — 보고서 & UI
# ═══════════════════════════════════════════════════════════
s8 = add_slide(); bg(s8)
slide_header(s8, "보고서 및 웹 인터페이스", "분석 결과 시각화 및 UI 구성", tag="06 UI & REPORT")

# 웹 UI 구성
box(s8, 0.4, 2.55, 5.9, 3.6, fill=BG_CARD)
accent_bar(s8, 0.4, 2.55, h=3.6, color=PRIMARY)
txt(s8, "웹 인터페이스 (localhost:8080)", 0.6, 2.65, 5.5, 0.4, size=13, bold=True, color=PRIMARY)
ui_items = [
    "• 메인 페이지 — .onion 도메인 입력 폼",
    "• 최근 분석 기록 (localStorage 저장)",
    "• 4가지 기능 소개 카드",
    "• 분석 중 로딩 스피너",
    "• 결과 페이지 — iframe 보고서 렌더링",
    "• 보고서 다운로드 버튼",
]
for j, item in enumerate(ui_items):
    txt(s8, item, 0.6, 3.15+j*0.48, 5.5, 0.38, size=12, color=LIGHT)

# 보고서 구성
box(s8, 6.55, 2.55, 6.3, 3.6, fill=BG_CARD)
accent_bar(s8, 6.55, 2.55, h=3.6, color=AMBER)
txt(s8, "HTML 분석 보고서 구성", 6.75, 2.65, 6.0, 0.4, size=13, bold=True, color=AMBER)
report_items = [
    ("요약 카드",         "접근성 · 위험도 · 유형 · CoDA 한눈에"),
    ("AI 사이트 분석",    "GPT-4o-mini 분석 결과 (API 키 필요)"),
    ("접근성 정보",       "HTTP 상태코드 · HTML 수집 여부"),
    ("검색 색인 정보",    "Ahmia · DuckDuckGo 노출 현황"),
    ("사이트 유형 분류",  "BART 분류 + 한국어 카테고리 차트"),
    ("CoDA 범죄 분류",    "DarkBERT 분류 + 확률 분포 차트"),
]
for j, (name, desc) in enumerate(report_items):
    t = 3.15 + j * 0.48
    txt(s8, f"• {name}", 6.75, t, 2.3, 0.3, size=12, bold=True, color=LIGHT)
    txt(s8, desc, 8.9, t, 4.0, 0.3, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — 기술 스택
# ═══════════════════════════════════════════════════════════
s9 = add_slide(); bg(s9)
slide_header(s9, "기술 스택", "사용된 프레임워크, 라이브러리, AI 모델", tag="07 TECH STACK")

stacks = [
    ("AI / ML 모델", [
        ("DarkBERT", "s2w-ai/darkbert", "다크웹 특화 BERT 모델"),
        ("BART",     "facebook/bart-large-mnli", "Zero-shot 분류"),
        ("GPT-4o-mini", "OpenAI API", "자연어 위험도 분석"),
        ("LogisticRegression", "scikit-learn", "CoDA 최종 분류기"),
    ], PRIMARY),
    ("백엔드", [
        ("Flask",    "Python 3.11", "웹 서버 / REST API"),
        ("Requests", "SOCKS5 지원", "Tor 경유 HTTP 요청"),
        ("Stem",     "Tor Control", "Tor 회로 제어"),
        ("BeautifulSoup4", "HTML 파싱", "콘텐츠 추출"),
    ], CYAN),
    ("프론트엔드 / 기타", [
        ("Pretendard", "웹폰트", "UI 타이포그래피"),
        ("Matplotlib", "차트 생성", "분석 결과 시각화"),
        ("python-pptx", "PPT 생성", "보고서 자동화"),
        ("Tor Browser", "SOCKS5:9050", "다크웹 접속 프록시"),
    ], GREEN),
]

for i, (cat, items, color) in enumerate(stacks):
    l = 0.4 + i * 4.3
    box(s9, l, 2.55, 4.0, 3.9, fill=BG_CARD)
    cb = s9.shapes.add_shape(1, Inches(l), Inches(2.55), Inches(4.0), Inches(0.06))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s9, cat, l+0.15, 2.65, 3.7, 0.4, size=13, bold=True, color=color)
    for j, (name, sub1, sub2) in enumerate(items):
        t = 3.15 + j * 0.82
        txt(s9, name, l+0.15, t, 3.7, 0.3, size=12, bold=True, color=WHITE)
        txt(s9, f"{sub1}  ·  {sub2}", l+0.15, t+0.3, 3.7, 0.28, size=10, color=MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — 마무리
# ═══════════════════════════════════════════════════════════
s10 = add_slide(); bg(s10)

txt(s10, "DARKWEB INTELLIGENCE", 1.5, 2.0, 10, 0.4, size=11, bold=True,
    color=PRIMARY, align=PP_ALIGN.CENTER)

title_box2 = s10.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.2))
title_box2.word_wrap = False
tf2 = title_box2.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "감사합니다"
run2.font.size = Pt(52)
run2.font.bold = True
run2.font.color.rgb = WHITE
run2.font.name = "Pretendard"

txt(s10, "Darkweb Domain Analyzer", 1.5, 3.75, 10, 0.4, size=16,
    color=INDIGO_SOFT, align=PP_ALIGN.CENTER)

d2 = s10.shapes.add_shape(1, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.03))
d2.fill.solid(); d2.fill.fore_color.rgb = PRIMARY; d2.line.fill.background()

txt(s10, "GitHub: github.com/junseong127/darkweb-analyzer", 1.5, 4.55, 10, 0.35,
    size=12, color=MUTED, align=PP_ALIGN.CENTER)
txt(s10, "보안 연구 목적으로 개발  |  S2W CoDA 데이터셋 활용", 1.5, 4.95, 10, 0.35,
    size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ── 저장 ───────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "Darkweb_Analyzer_발표.pptx")
prs.save(out)
print(f"✅ PPT 생성 완료: {out}")
